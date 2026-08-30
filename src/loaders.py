from __future__ import annotations

import hashlib
import io
import json
import tempfile
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Iterable, List, Optional

import fitz  # PyMuPDF
import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from langchain_core.documents import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .vision import GroqVisionExtractor, guess_mime_type

SUPPORTED_EXTENSIONS = {
    ".pdf", ".txt", ".md", ".markdown", ".docx", ".pptx", ".xlsx", ".xlsm",
    ".csv", ".tsv", ".html", ".htm", ".json", ".xml", ".yaml", ".yml",
    ".eml", ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".xml", ".yaml", ".yml"}


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _meta(source: str, file_hash: str, content_type: str, locator: str = "", **extra) -> dict:
    data = {
        "source_file": source,
        "source": source,
        "file_hash": file_hash,
        "content_type": content_type,
        "locator": locator,
    }
    data.update(extra)
    return data


def _doc(text: str, source: str, file_hash: str, content_type: str = "text", locator: str = "", **extra) -> List[Document]:
    text = (text or "").strip()
    if not text:
        return []
    return [Document(page_content=text, metadata=_meta(source, file_hash, content_type, locator, **extra))]


def _markdown_table(rows) -> str:
    cleaned = [["" if v is None else str(v).strip() for v in row] for row in rows if row]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    cleaned = [row + [""] * (width - len(row)) for row in cleaned]
    header, body = cleaned[0], cleaned[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _vision_call(vision: Optional[GroqVisionExtractor], data: bytes, mime: str, instruction: str) -> str:
    if not vision:
        return ""
    try:
        return vision.describe_image_bytes(data, mime, instruction)
    except Exception as exc:
        return f"[Visual extraction failed: {type(exc).__name__}]"


def _shape_position(shape) -> str:
    """Return safe PPTX shape coordinates without assuming they are always populated."""
    left = getattr(shape, "left", None)
    top = getattr(shape, "top", None)
    if left is None or top is None:
        return "position unavailable"
    try:
        return f"x={int(left)}, y={int(top)}"
    except (TypeError, ValueError, OverflowError):
        return "position unavailable"


def _load_pdf(path: Path, source: str, file_hash: str, vision: Optional[GroqVisionExtractor], max_vision_items: int) -> List[Document]:
    pdf = fitz.open(str(path))
    docs: List[Document] = []
    vision_used = 0
    try:
        for idx, page in enumerate(pdf):
            page_no = idx + 1
            parts: List[str] = []
            text = page.get_text("text").strip()
            if text:
                parts.append("TEXT\n" + text)

            try:
                tables = page.find_tables().tables
            except Exception:
                tables = []
            for table_no, table in enumerate(tables, start=1):
                md = _markdown_table(table.extract())
                if md:
                    parts.append(f"TABLE {table_no}\n{md}")

            try:
                image_count = len(page.get_images(full=True))
            except Exception:
                image_count = 0
            try:
                drawing_count = len(page.get_drawings())
            except Exception:
                drawing_count = 0

            needs_vision = bool(vision and (image_count > 0 or drawing_count >= 4 or len(text) < 80))
            if needs_vision and (max_vision_items == 0 or vision_used < max_vision_items):
                vision_used += 1
                pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
                visual = _vision_call(
                    vision,
                    pix.tobytes("png"),
                    "image/png",
                    f"Page {page_no} of {source}. Extract scanned text, tables, charts, diagrams, formulas, labels and relationships not reliably captured by normal parsing.",
                )
                if visual:
                    parts.append("VISUAL/LAYOUT EVIDENCE\n" + visual)

            if parts:
                docs.append(
                    Document(
                        page_content="\n\n".join(parts),
                        metadata=_meta(
                            source, file_hash, "page", f"page {page_no}",
                            page=idx, page_number=page_no, total_pages=len(pdf),
                        ),
                    )
                )
    finally:
        pdf.close()
    return docs


def _load_docx(path: Path, source: str, file_hash: str, vision: Optional[GroqVisionExtractor], max_vision_items: int) -> List[Document]:
    docx = DocxDocument(str(path))
    docs: List[Document] = []
    paragraphs = [p.text.strip() for p in docx.paragraphs if p.text.strip()]
    if paragraphs:
        docs += _doc("\n\n".join(paragraphs), source, file_hash, "text", "document body")

    for i, table in enumerate(docx.tables, start=1):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        md = _markdown_table(rows)
        if md:
            docs += _doc(f"TABLE {i}\n{md}", source, file_hash, "table", f"table {i}")

    if vision:
        used = 0
        seen = set()
        for rel in docx.part.rels.values():
            target = getattr(rel, "target_part", None)
            content_type = getattr(target, "content_type", "") if target else ""
            if not content_type.startswith("image/"):
                continue
            blob = getattr(target, "blob", b"")
            digest = hashlib.sha256(blob).hexdigest() if blob else ""
            if not blob or digest in seen:
                continue
            seen.add(digest)
            if max_vision_items and used >= max_vision_items:
                break
            used += 1
            visual = _vision_call(vision, blob, content_type, f"Embedded image {used} in {source}. Extract all text, chart, diagram or visual evidence.")
            if visual:
                docs += _doc(visual, source, file_hash, "visual", f"embedded image {used}")
    return docs


def _ppt_chart_text(shape, chart_index: int) -> str:
    chart = shape.chart
    lines = [f"CHART {chart_index}"]
    try:
        if chart.has_title:
            lines.append("Title: " + chart.chart_title.text_frame.text.strip())
    except Exception:
        pass
    try:
        categories = [str(x) for x in chart.plots[0].categories]
        if categories:
            lines.append("Categories: " + ", ".join(categories))
    except Exception:
        pass
    for i, series in enumerate(getattr(chart, "series", []) or [], start=1):
        try:
            lines.append(f"Series {i} ({series.name}): " + ", ".join(str(v) for v in series.values))
        except Exception:
            lines.append(f"Series {i}: [values unavailable]")
    return "\n".join(lines)


def _load_pptx(path: Path, source: str, file_hash: str, vision: Optional[GroqVisionExtractor], max_vision_items: int) -> List[Document]:
    prs = Presentation(str(path))
    docs: List[Document] = []
    vision_used = 0
    for slide_no, slide in enumerate(prs.slides, start=1):
        blocks: List[str] = []
        spatial: List[str] = []
        table_no = chart_no = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    blocks.append(text)
                    spatial.append(f"Text at {_shape_position(shape)}: {text}")
            if getattr(shape, "has_table", False):
                table_no += 1
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                md = _markdown_table(rows)
                if md:
                    blocks.append(f"TABLE {table_no}\n{md}")
            if getattr(shape, "has_chart", False):
                chart_no += 1
                blocks.append(_ppt_chart_text(shape, chart_no))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and vision:
                if max_vision_items == 0 or vision_used < max_vision_items:
                    vision_used += 1
                    visual = _vision_call(
                        vision, shape.image.blob, shape.image.content_type,
                        f"Picture on slide {slide_no} of {source}. Extract text, chart, diagram, labels, arrows and relationships.",
                    )
                    if visual:
                        blocks.append("VISUAL\n" + visual)
        if len(spatial) > 1:
            blocks.append("SPATIAL TEXT ELEMENTS\n" + "\n".join(spatial))
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append("SPEAKER NOTES\n" + notes)
        except Exception:
            pass
        if blocks:
            docs.append(Document(page_content=f"SLIDE {slide_no}\n\n" + "\n\n".join(blocks), metadata=_meta(source, file_hash, "slide", f"slide {slide_no}", slide_number=slide_no)))
    return docs


def _load_excel(path: Path, source: str, file_hash: str, vision: Optional[GroqVisionExtractor], max_vision_items: int) -> List[Document]:
    docs: List[Document] = []
    sheets = pd.read_excel(path, sheet_name=None, engine="openpyxl")
    wb_formula = load_workbook(path, data_only=False)
    wb_values = load_workbook(path, data_only=True)
    vision_used = 0

    for sheet_name, df in sheets.items():
        if df.empty:
            docs += _doc(f"SHEET: {sheet_name}\n[Empty sheet]", source, file_hash, "table", f"sheet {sheet_name}")
        else:
            batch = 120
            for start in range(0, len(df), batch):
                part = df.iloc[start:start + batch]
                locator = f"sheet {sheet_name}, rows {start + 1}-{start + len(part)}"
                docs += _doc(f"SHEET: {sheet_name}\n{part.to_markdown(index=False)}", source, file_hash, "table", locator)

        formulas = []
        for row in wb_formula[sheet_name].iter_rows():
            for cell in row:
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    formulas.append(f"{cell.coordinate}: {cell.value}")
        if formulas:
            docs += _doc(f"SHEET: {sheet_name}\nFORMULAS\n" + "\n".join(formulas), source, file_hash, "formula", f"sheet {sheet_name}, formulas")

        ws = wb_values[sheet_name]
        for chart_no, chart in enumerate(getattr(ws, "_charts", []) or [], start=1):
            lines = [f"SHEET: {sheet_name}", f"CHART {chart_no} ({chart.__class__.__name__})"]
            for series_no, series in enumerate(getattr(chart, "ser", []) or [], start=1):
                lines.append(f"Series {series_no}: chart series present in workbook")
            docs += _doc("\n".join(lines), source, file_hash, "chart", f"sheet {sheet_name}, chart {chart_no}")

        if vision:
            for image_no, image in enumerate(getattr(ws, "_images", []) or [], start=1):
                if max_vision_items and vision_used >= max_vision_items:
                    break
                vision_used += 1
                try:
                    blob = image._data()
                except Exception:
                    continue
                visual = _vision_call(vision, blob, "image/png", f"Embedded image {image_no} on Excel sheet {sheet_name} in {source}.")
                if visual:
                    docs += _doc(visual, source, file_hash, "visual", f"sheet {sheet_name}, image {image_no}")
    return docs


def _load_html(path: Path, source: str, file_hash: str) -> List[Document]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    docs = _doc("\n".join(x.strip() for x in soup.get_text("\n").splitlines() if x.strip()), source, file_hash, "html", "document body")
    try:
        tables = pd.read_html(io.StringIO(raw))
    except Exception:
        tables = []
    for i, df in enumerate(tables, start=1):
        docs += _doc(f"HTML TABLE {i}\n{df.to_markdown(index=False)}", source, file_hash, "table", f"HTML table {i}")
    return docs


def _load_delimited(path: Path, source: str, file_hash: str, sep: str) -> List[Document]:
    df = pd.read_csv(path, sep=sep)
    docs: List[Document] = []
    batch = 150
    if df.empty:
        return _doc("TABLE\n[Empty table]", source, file_hash, "table", "table")
    for start in range(0, len(df), batch):
        part = df.iloc[start:start + batch]
        docs += _doc("TABLE\n" + part.to_markdown(index=False), source, file_hash, "table", f"rows {start + 1}-{start + len(part)}")
    return docs


def _load_eml(path: Path, source: str, file_hash: str) -> List[Document]:
    msg = BytesParser(policy=policy.default).parsebytes(path.read_bytes())
    header = [f"From: {msg.get('from', '')}", f"To: {msg.get('to', '')}", f"Subject: {msg.get('subject', '')}", f"Date: {msg.get('date', '')}"]
    bodies = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == "text/plain":
                try:
                    bodies.append(part.get_content())
                except Exception:
                    pass
    else:
        try:
            bodies.append(msg.get_content())
        except Exception:
            pass
    return _doc("\n".join(header) + "\n\n" + "\n\n".join(bodies), source, file_hash, "email", "email body")


def _load_path(path: Path, source: str, file_hash: str, vision: Optional[GroqVisionExtractor], max_vision_items: int) -> List[Document]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path, source, file_hash, vision, max_vision_items)
    if suffix in TEXT_EXTENSIONS:
        return _doc(path.read_text(encoding="utf-8", errors="replace"), source, file_hash, "text", "document body")
    if suffix in {".html", ".htm"}:
        return _load_html(path, source, file_hash)
    if suffix == ".json":
        try:
            obj = json.loads(path.read_text(encoding="utf-8", errors="replace"))
            text = json.dumps(obj, indent=2, ensure_ascii=False)
        except Exception:
            text = path.read_text(encoding="utf-8", errors="replace")
        return _doc(text, source, file_hash, "json", "document body")
    if suffix == ".csv":
        return _load_delimited(path, source, file_hash, ",")
    if suffix == ".tsv":
        return _load_delimited(path, source, file_hash, "\t")
    if suffix in {".xlsx", ".xlsm"}:
        return _load_excel(path, source, file_hash, vision, max_vision_items)
    if suffix == ".docx":
        return _load_docx(path, source, file_hash, vision, max_vision_items)
    if suffix == ".pptx":
        return _load_pptx(path, source, file_hash, vision, max_vision_items)
    if suffix == ".eml":
        return _load_eml(path, source, file_hash)
    if suffix in IMAGE_EXTENSIONS:
        if not vision:
            raise ValueError(f"{source} is image-only. Enable multimodal extraction and provide GROQ_API_KEY.")
        visual = _vision_call(vision, path.read_bytes(), guess_mime_type(source), f"Standalone document image: {source}.")
        return _doc(visual, source, file_hash, "visual", "image")
    raise ValueError(f"Unsupported file type: {suffix}")


def load_uploaded_files(
    uploaded_files: Iterable,
    *,
    api_key: str = "",
    vision_model: str = "",
    enable_vision: bool = True,
    max_vision_items: int = 20,
) -> List[Document]:
    """Load Streamlit UploadedFile-like objects into LangChain Documents."""
    vision = GroqVisionExtractor(api_key=api_key, model_name=vision_model) if enable_vision and api_key and vision_model else None
    all_docs: List[Document] = []
    for uploaded in uploaded_files:
        name = Path(uploaded.name).name
        suffix = Path(name).suffix.lower()
        if suffix not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"{name}: unsupported format. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}")
        data = uploaded.getvalue()
        file_hash = sha256_bytes(data)
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(data)
            tmp_path = Path(tmp.name)
        try:
            try:
                extracted = _load_path(tmp_path, name, file_hash, vision, max_vision_items)
            except Exception as exc:
                raise ValueError(f"{name}: {type(exc).__name__}: {exc}") from exc
            all_docs.extend(extracted)
        finally:
            tmp_path.unlink(missing_ok=True)
    return all_docs


def load_paths(
    paths: Iterable[str | Path],
    *,
    api_key: str = "",
    vision_model: str = "",
    enable_vision: bool = True,
    max_vision_items: int = 20,
) -> List[Document]:
    """Load local files outside Streamlit."""
    vision = GroqVisionExtractor(api_key=api_key, model_name=vision_model) if enable_vision and api_key and vision_model else None
    docs: List[Document] = []
    for raw in paths:
        path = Path(raw)
        if path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {path.suffix.lower()}")
        data = path.read_bytes()
        docs.extend(_load_path(path, path.name, sha256_bytes(data), vision, max_vision_items))
    return docs
